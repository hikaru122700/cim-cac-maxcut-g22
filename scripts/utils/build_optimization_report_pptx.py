"""BrainPad『数理最適化研修』(105p) を、研究室発表用に 20 枚に再構成して .pptx を生成する。

研究室の文脈(MaxCut/CIM/CAC/Optuna)に接続するため、終盤に「研究との橋渡し」を 3 枚追加。
出力: docs/optimization_training_report.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ============================================================
# テーマ
# ============================================================
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)   # 濃紺
COLOR_ACCENT = RGBColor(0xC0, 0x39, 0x2B)    # 赤
COLOR_SUB = RGBColor(0x4A, 0x4A, 0x4A)       # ダークグレー
COLOR_LIGHT = RGBColor(0x95, 0xA5, 0xA6)     # ライトグレー
COLOR_HIGHLIGHT = RGBColor(0xF1, 0xC4, 0x0F) # 黄

FONT_JA = "Yu Gothic"
FONT_EN = "Consolas"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ============================================================
# 補助関数
# ============================================================
def add_title_bar(slide, title: str, section: str | None = None):
    """画面上部に色付き帯 + タイトル。"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W_IN), Inches(0.85),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(11.5), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_JA
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if section:
        tb2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.55), Inches(11.5), Inches(0.3))
        p2 = tb2.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = section
        run2.font.name = FONT_JA
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(0xD0, 0xE2, 0xF2)


def add_footer(slide, idx: int, total: int):
    """フッターにページ番号と出典。"""
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = (
        f"BrainPad 数理最適化研修 報告 (五十嵐研)  /  "
        f"出典: 数理最適化研修 (BrainPad Inc., 2025)  /  {idx} / {total}"
    )
    run.font.name = FONT_JA
    run.font.size = Pt(9)
    run.font.color.rgb = COLOR_LIGHT


def add_textbox(
    slide, left, top, width, height, text: str,
    size: int = 14, bold: bool = False,
    color: RGBColor = COLOR_SUB, align=PP_ALIGN.LEFT,
    font_name: str = FONT_JA,
):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide, left, top, width, height,
    bullets: list[str | tuple[str, int]],
    size: int = 14, leading_indent: int = 0,
    color: RGBColor = COLOR_SUB,
):
    """箇条書きを 1 つの textbox にまとめる。

    bullets の要素が tuple のときは (text, level) を表す。
    """
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = level + leading_indent
        bullet = "•" if level == 0 else "‣" if level == 1 else "–"
        run = p.add_run()
        run.text = f"{bullet}  {text}"
        run.font.name = FONT_JA
        run.font.size = Pt(size - level * 1)
        run.font.color.rgb = color
    return tb


def add_box(
    slide, left, top, width, height,
    title: str, body: list[str],
    fill: RGBColor = RGBColor(0xEC, 0xF2, 0xF9),
    title_color: RGBColor = COLOR_PRIMARY,
):
    """色付きパネル(枠 + タイトル + 箇条書き)。"""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = title_color
    rect.line.width = Pt(1.2)
    rect.shadow.inherit = False

    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.08),
        Inches(width - 0.3), Inches(0.4),
    )
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_JA
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = title_color

    add_bullets(
        slide, left + 0.2, top + 0.5,
        width - 0.4, height - 0.6,
        body, size=12, color=COLOR_SUB,
    )


def add_section_divider(slide, section_no: str, title: str):
    """章扉スライド。"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    add_textbox(
        slide, 0.5, 2.7, 12.5, 1.2,
        section_no, size=64, bold=True,
        color=RGBColor(0xF1, 0xC4, 0x0F),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, 0.5, 3.9, 12.5, 1.2,
        title, size=36, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )


# ============================================================
# スライド本体
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    slides_meta: list[str] = []

    # =======================================
    # Slide 1: 表紙
    # =======================================
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2),
        Inches(SLIDE_W_IN), Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_HIGHLIGHT
    accent.line.fill.background()
    add_textbox(s, 0.5, 1.3, 12.5, 1.0, "数理最適化研修 報告",
                size=46, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 2.3, 12.5, 0.6,
                "── LP / IP / MIP を学んで MaxCut の解き方を整理する ──",
                size=20, color=RGBColor(0xD0, 0xE2, 0xF2),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 3.6, 12.5, 0.5,
                "原典: BrainPad『数理最適化研修』 全 105 ページ",
                size=14, color=RGBColor(0xD0, 0xE2, 0xF2),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 4.2, 12.5, 0.5,
                "発表者: 五十嵐研  CIM-CAC-MaxCut チーム",
                size=16, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.5, 6.6, 12.5, 0.5, "2026-05",
                size=12, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    slides_meta.append("表紙")

    # =======================================
    # Slide 2: 本日のアジェンダ
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "本日のアジェンダ", "20 枚で 105 ページを縮約")
    add_bullets(s, 0.6, 1.2, 6.0, 5.5, [
        ("Part I  数理最適化とは — 第 1〜4 枚", 0),
        ("最適化問題の 3 要素と分類", 1),
        ("解くまでの流れ・ソルバーの位置づけ", 1),
        ("Part II  LP / IP / MIP — 第 5〜13 枚", 0),
        ("線形計画 (LP): 飲料生産問題", 1),
        ("整数計画 (IP, 0-1): ナップサック", 1),
        ("混合整数計画 (MIP): 施設配置・チーム分け", 1),
        ("Part III  研究室との接続 — 第 14〜18 枚", 0),
        ("MaxCut を IP として書く", 1),
        ("MILP ソルバー vs CIM/CAC の住み分け", 1),
        ("Optuna も「最適化問題」だった", 1),
        ("まとめ・参考文献 — 第 19〜20 枚", 0),
    ], size=14)
    add_box(s, 7.0, 1.2, 5.8, 5.5,
            "本資料のゴール",
            [
                "BrainPad の研修内容(LP/IP/MIP)を、研究室で",
                "自分たちが解いている MaxCut / CIM とどう繋がるか",
                "という観点に置き換えて整理する。",
                "",
                "テクニックそのものではなく、",
                "「自分たちの問題を数理最適化の言葉で",
                "  記述・分類する力」を持ち帰る。",
                "",
                "→ MaxCut は 0-1 IP の典型問題",
                "→ CIM はフルスクラッチ系アルゴリズム",
                "→ Optuna は連続ブラックボックス最適化",
            ],
            fill=RGBColor(0xFD, 0xF6, 0xE3),
            title_color=COLOR_ACCENT)
    add_footer(s, 2, 20)
    slides_meta.append("アジェンダ")

    # =======================================
    # Slide 3: 数理最適化とは
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "数理最適化とは何か", "①数理最適化イントロダクション")
    add_textbox(s, 0.6, 1.1, 12.0, 0.6,
                "制約条件のもとで目的関数を最大化/最小化する解を求める「問題」と、それを通じた意思決定の「技術」",
                size=15, bold=True, color=COLOR_PRIMARY)

    # 3 要素のカード
    cards = [
        ("決定変数", "最適化で決めたい値\n  x ∈ ℝⁿ (連続) や {0,1}ⁿ (離散)",
         RGBColor(0xEC, 0xF2, 0xF9), COLOR_PRIMARY),
        ("目的関数 f(x)", "最大化 / 最小化したい指標\n  売上、損失、コスト、…",
         RGBColor(0xFD, 0xF6, 0xE3), COLOR_ACCENT),
        ("制約条件 g(x) ≤ 0", "実行可能領域を絞る式\n  容量、需要、論理関係、…",
         RGBColor(0xE8, 0xF8, 0xF5), RGBColor(0x16, 0xA0, 0x85)),
    ]
    for i, (t, b, fill, tc) in enumerate(cards):
        left = 0.6 + i * 4.2
        add_box(s, left, 2.0, 3.9, 2.0,
                t, [ln for ln in b.split("\n")],
                fill=fill, title_color=tc)

    # 機械学習との対比
    add_textbox(s, 0.6, 4.4, 12.0, 0.5,
                "機械学習(予測) ↔ 数理最適化(意思決定)",
                size=16, bold=True, color=COLOR_PRIMARY)
    add_bullets(s, 0.6, 4.95, 12.0, 1.8, [
        "ML: 「次にどうなる?」 を入力データから当てるモデル → 説明変数 ⇒ 目的変数",
        "MO: 「どうしたら良い?」 を目的関数の最大化で答える → 数理モデル ⇒ 行動",
        "本研究室の CIM 研究も「ML がやらないこと」を扱う立場 — 意思決定 (=どう切るか) に近い",
    ], size=13)
    add_footer(s, 3, 20)
    slides_meta.append("数理最適化とは")

    # =======================================
    # Slide 4: 最適化問題の分類
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "最適化問題の分類", "用語と地図")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "決定変数と目的・制約が「線形か / 連続か」で大別される",
                size=14, color=COLOR_SUB)

    # 4 象限
    quadrants = [
        ("連続  ×  線形", "LP\nLinear Programming",
         "実数決定変数 + 線形 f, g\n例: 生産量配分", RGBColor(0xEC, 0xF2, 0xF9)),
        ("連続  ×  非線形", "NLP / QP",
         "実数決定変数 + 非線形\n例: ML の損失関数最小化\n本研究室の CIM パラ調整",
         RGBColor(0xFD, 0xF6, 0xE3)),
        ("離散  ×  線形", "IP / MIP",
         "整数 (0/1) 決定変数 + 線形\n例: ナップサック、施設配置\n★MaxCut もここ",
         RGBColor(0xFA, 0xE5, 0xD3)),
        ("離散  ×  非線形", "INLP / QUBO",
         "整数決定変数 + 非線形\n例: QUBO, Ising モデル\n★CIM が解いているのは実質これ",
         RGBColor(0xE8, 0xF8, 0xF5)),
    ]
    for i, (label, name, body, fill) in enumerate(quadrants):
        r = i // 2; c = i % 2
        left = 0.6 + c * 6.2
        top = 1.7 + r * 2.65
        add_box(s, left, top, 6.0, 2.5, label,
                [name, "", body],
                fill=fill, title_color=COLOR_PRIMARY)

    add_textbox(s, 0.6, 6.95, 12.0, 0.35,
                "★ = 自分たちの研究テーマがここに位置することを示す",
                size=11, color=COLOR_ACCENT, bold=True)
    add_footer(s, 4, 20)
    slides_meta.append("分類")

    # =======================================
    # Slide 5: 最適化案件の進め方
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "最適化案件の進め方", "定式化 → 求解 → 可視化 → 改良")
    steps = [
        ("現実問題", "言語的に書かれた要件\n生産計画・配送・配置 …"),
        ("最適化問題", "目的関数 + 決定変数\n+ 制約式 に分解"),
        ("計算結果", "ソルバーや独自アルゴ\nが返す数値解"),
        ("問題解決", "可視化で違和感を検出\n要件と摺り合わせ"),
    ]
    for i, (t, b) in enumerate(steps):
        left = 0.4 + i * 3.25
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(2.0), Inches(3.0), Inches(2.6),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xEC, 0xF2, 0xF9)
        box.line.color.rgb = COLOR_PRIMARY
        add_textbox(s, left + 0.1, 2.1, 2.8, 0.5,
                    t, size=15, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(s, left + 0.1, 2.7, 2.8, 1.7,
                    b, size=11, color=COLOR_SUB, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(left + 3.05), Inches(3.0),
                Inches(0.2), Inches(0.6),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.fill.background()

    # 反復矢印
    add_textbox(s, 0.4, 4.85, 12.5, 0.5,
                "↑ 1 回でうまく行くことは稀。再定式化 / 再アルゴリズム選択を反復する",
                size=13, color=COLOR_ACCENT, align=PP_ALIGN.CENTER, bold=True)

    add_bullets(s, 0.6, 5.5, 12.0, 1.5, [
        "研究室文脈の翻訳:",
        ("「現実問題」=最大カット数を取りたい / 振幅が綺麗に二極化したい", 1),
        ("「定式化」 =MaxCut → 0-1 IP / Ising → SDE", 1),
        ("「ソルバー」=Cbc / Gurobi / SCIP  vs  CIM / CAC / SA / SB", 1),
        ("「可視化・改良」=cut 分布ヒスト・振幅軌跡 → パラメータ再調整", 1),
    ], size=12)
    add_footer(s, 5, 20)
    slides_meta.append("案件の進め方")

    # =======================================
    # Slide 6: 解くアプローチの 2 分類
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "汎用ソルバー vs フルスクラッチ", "性能と汎用性のトレードオフ")
    add_box(s, 0.6, 1.2, 6.0, 5.6,
            "(1) 汎用ソルバー",
            [
                "予めアルゴリズムを実装したパッケージ",
                "Gurobi (有償) / SCIP / Cbc / CP-SAT",
                "",
                "○ 制約変更に強い・保守容易",
                "○ 「裏で何をやっているか」を意識せず使える",
                "× 問題固有の構造は使えない",
                "× 大規模になると現実時間で解けない",
                "",
                "案件ではこれが第一選択",
                "(BrainPad 研修も PuLP + Cbc 想定)",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_box(s, 6.8, 1.2, 6.0, 5.6,
            "(2) フルスクラッチ・問題固有アルゴリズム",
            [
                "自分でアルゴリズムを 1 から実装",
                "ヒューリスティクス / メタヒューリスティクス",
                "",
                "○ 問題固有の性質を最大活用 → 高速・大規模対応",
                "× 保守性が低い・汎用性なし",
                "× 設計・チューニング・検証コストが高い",
                "",
                "★研究室の CIM / CAC / SB / SA はここに該当",
                "  汎用 MILP ソルバーが大規模 G-set を解けない領域で、",
                "  Ising 物理を使って高速に近似解を得る",
            ],
            fill=RGBColor(0xFD, 0xF6, 0xE3), title_color=COLOR_ACCENT)
    add_footer(s, 6, 20)
    slides_meta.append("ソルバー2分類")

    # =======================================
    # Slide 7: 章扉 II
    # =======================================
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "Part II", "LP / IP / MIP の定式化")
    slides_meta.append("Part II 扉")

    # =======================================
    # Slide 8: LP - 線形計画問題
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "線形計画問題 (LP)", "② 線形計画 定式化")
    add_textbox(s, 0.6, 1.1, 12.0, 0.6,
                "いくつかの一次不等式/等式を制約として、ある一次関数を最大化・最小化する問題",
                size=15, bold=True, color=COLOR_PRIMARY)
    add_box(s, 0.6, 1.9, 6.0, 4.8,
            "基本性質",
            [
                "最大化と最小化は等価",
                "  max f(x)  ⇔  min -f(x)",
                "",
                "不等式制約と等式制約も等価",
                "  ax ≤ b  ⇔  ax + s = b, s ≥ 0",
                "",
                "実行可能領域は多面体になる",
                "最適解は通常その「頂点」に位置 (単体法の根拠)",
                "",
                "ソルバーは多項式時間で解ける",
                "  → LP は事実上「解ける問題」",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_box(s, 6.8, 1.9, 6.0, 4.8,
            "一般形",
            [
                "minimize    cᵀ x",
                "subject to  A x  ≤  b",
                "            x  ≥  0",
                "",
                "  c ∈ ℝⁿ        目的関数の重み",
                "  A ∈ ℝᵐˣⁿ      制約行列",
                "  b ∈ ℝᵐ        制約の右辺",
                "  x ∈ ℝⁿ        決定変数 (実数!)",
                "",
                "決定変数が実数なのが LP の本質",
                "整数化されると一気に難しくなる ⇒ 次節 IP",
            ],
            fill=RGBColor(0xE8, 0xF8, 0xF5), title_color=RGBColor(0x16, 0xA0, 0x85))
    add_footer(s, 8, 20)
    slides_meta.append("LP 定義")

    # =======================================
    # Slide 9: LP 例題 飲料生産
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "LP 例題: 飲料の生産量最適化", "定式化の 3 ステップを体験")
    add_box(s, 0.6, 1.1, 6.0, 5.6,
            "問題設定",
            [
                "売値: フルーツジュース 3,000 円/kg, ティー 2,000 円/kg",
                "材料: リンゴ・オレンジ(在庫 500 / 300)",
                "1 kg 製造あたり必要量:",
                "  ジュース → リンゴ 2, オレンジ 4",
                "  ティー   → リンゴ 5, オレンジ 2",
                "売上を最大化するには何 kg ずつ作る?",
                "",
                "「現実問題」をどうやって",
                "  決定変数 / 目的関数 / 制約 に分解するか",
                "がこの章で持ち帰るべきスキル",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_box(s, 6.8, 1.1, 6.0, 5.6,
            "定式化",
            [
                "決定変数: x₁ = ジュース量, x₂ = ティー量 (kg, ≥ 0)",
                "",
                "目的関数:",
                "  maximize   3000 x₁ + 2000 x₂",
                "",
                "制約条件:",
                "  リンゴ:    2 x₁ + 5 x₂  ≤ 500",
                "  オレンジ:  4 x₁ + 2 x₂  ≤ 300",
                "  非負:      x₁, x₂ ≥ 0",
                "",
                "解: x₁ = 31.25, x₂ = 87.5,",
                "  売上 = 268,750 円",
            ],
            fill=RGBColor(0xFD, 0xF6, 0xE3), title_color=COLOR_ACCENT)
    add_footer(s, 9, 20)
    slides_meta.append("LP 例題")

    # =======================================
    # Slide 10: LP 実装 PuLP
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "PuLP + Cbc による実装", "③ 線形計画 実装")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "モデラー (PuLP) で定式化を Python で書き、ソルバー (Cbc) が裏で解く",
                size=13, color=COLOR_SUB)

    code = (
        "import pulp\n"
        "\n"
        "model = pulp.LpProblem(name='飲料生産', sense=pulp.LpMaximize)\n"
        "\n"
        "x1 = pulp.LpVariable('x1', lowBound=0, cat='Continuous')\n"
        "x2 = pulp.LpVariable('x2', lowBound=0, cat='Continuous')\n"
        "\n"
        "model += 3000 * x1 + 2000 * x2          # 目的関数\n"
        "model += 2 * x1 + 5 * x2 <= 500         # リンゴ制約\n"
        "model += 4 * x1 + 2 * x2 <= 300         # オレンジ制約\n"
        "\n"
        "model.solve()\n"
        "print('Status:', pulp.LpStatus[model.status])\n"
        "print('x1 =', x1.value(), 'x2 =', x2.value())\n"
        "# → Status: Optimal,  x1 = 31.25, x2 = 87.5"
    )
    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.7),
                              Inches(7.5), Inches(5.0))
    rect = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(1.65), Inches(7.6), Inches(5.05),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    rect.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.75),
                              Inches(7.4), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_EN
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0xF8, 0xF8, 0xF2)

    add_box(s, 8.4, 1.7, 4.4, 5.0,
            "モデラー作成 6 ステップ",
            [
                "1. 問題定義",
                "    LpProblem(name, sense)",
                "2. 決定変数",
                "    LpVariable(name, bounds, cat)",
                "3. 目的関数",
                "    model += expr",
                "4. 制約式",
                "    model += expr <= rhs",
                "5. 求解",
                "    model.solve()",
                "6. 値取り出し",
                "    x.value()",
                "",
                "ML の fit→predict と相似形",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_footer(s, 10, 20)
    slides_meta.append("PuLP 実装")

    # =======================================
    # Slide 11: IP / 0-1 IP
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "整数計画 (IP) と 0-1 整数計画", "④ 整数計画問題")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "決定変数が整数値のみ。特に {0, 1} に限定されたものを 0-1 IP と呼ぶ",
                size=14, color=COLOR_SUB)
    add_box(s, 0.6, 1.7, 6.0, 5.0,
            "種類と難しさ",
            [
                "0-1 IP: 選択する/しない (binary)",
                "  → ナップサック、施設配置の「契約 or NOT」",
                "",
                "線形 IP: 整数 + 線形目的・制約",
                "  → ビンパッキングなど",
                "",
                "非線形 IP: 二乗誤差や log を含む",
                "  → PuLP + Cbc では解けない",
                "",
                "★MaxCut / Ising は二次形の 0-1 IP",
                "  決定変数 xᵢ ∈ {0, 1} が「集合 A か B か」",
                "  目的 = Σ (xᵢ - xⱼ)² over edges (i, j)",
            ],
            fill=RGBColor(0xFA, 0xE5, 0xD3), title_color=COLOR_ACCENT)
    add_box(s, 6.8, 1.7, 6.0, 5.0,
            "ナップサック問題の定式化",
            [
                "決定変数: xᵢ ∈ {0, 1}  商品 i を入れるか",
                "",
                "目的:  maximize  Σᵢ vᵢ xᵢ",
                "",
                "制約:  Σᵢ wᵢ xᵢ  ≤  C   (容量)",
                "",
                "実装は LP とほぼ同じ:",
                "  pulp.LpVariable.dicts('x', I, cat='Binary')",
                "",
                "ここで集合 (I) と総和記号 (Σ) で書く",
                "  → 入力サイズに依らない汎用コード",
                "",
                "MaxCut も「集合 + Σ + 0-1 変数」で全部書ける",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_footer(s, 11, 20)
    slides_meta.append("IP")

    # =======================================
    # Slide 12: MIP 施設配置
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "混合整数計画 (MIP) — 施設配置", "⑤ MIP 定式化")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "実数決定変数と整数決定変数が同居 — 現実問題の多くがここに該当",
                size=14, color=COLOR_SUB)
    add_box(s, 0.6, 1.7, 6.0, 5.0,
            "問題: 倉庫契約と輸送量",
            [
                "5 つの倉庫から 4 店舗へ製品を配送する",
                "  yᵢ ∈ {0, 1}      倉庫 i を契約するか        (整数!)",
                "  xᵢⱼ ∈ ℝ ≥ 0      i から j への輸送量          (実数!)",
                "",
                "目的:  min Σᵢ fᵢ yᵢ  +  Σᵢⱼ cᵢⱼ xᵢⱼ",
                "  (固定費 + 輸送費)",
                "",
                "制約:",
                "  Σⱼ xᵢⱼ  ≤  Mᵢ yᵢ        容量(big-M トリック)",
                "  Σᵢ xᵢⱼ  =  dⱼ            店舗 j の需要充足",
                "",
                "→ 整数 + 実数の両方を扱う = MIP",
            ],
            fill=RGBColor(0xFA, 0xE5, 0xD3), title_color=COLOR_ACCENT)
    add_box(s, 6.8, 1.7, 6.0, 5.0,
            "big-M トリックの読み方",
            [
                "  Σⱼ xᵢⱼ  ≤  Mᵢ yᵢ",
                "",
                "yᵢ = 0 のとき右辺 = 0 → 倉庫 i は 1 単位も出せない",
                "yᵢ = 1 のとき右辺 = Mᵢ → 倉庫 i の容量分だけ出せる",
                "",
                "= 「整数変数で連続変数のオン/オフを切る」",
                "",
                "★Ising / MaxCut でも同じ構造が随所に出る",
                "  σᵢ = ±1 で結合項 Jᵢⱼ σᵢ σⱼ をオン/オフ",
                "  「離散変数で連続物理量の挙動を切替える」",
                "",
                "MIP の感覚を持っておくと",
                "  自分の問題を MILP ソルバーに投げる選択肢が生まれる",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_footer(s, 12, 20)
    slides_meta.append("MIP 施設配置")

    # =======================================
    # Slide 13: ハード制約をソフトに
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "ハード制約 → ソフト制約", "実行可能性が壊れたときの逃がし方")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "等式 Σ = AvgPoints のように厳密に守らせると実行不能になりがち",
                size=14, color=COLOR_SUB)
    add_box(s, 0.6, 1.7, 6.0, 5.0,
            "ハード制約だと壊れる例",
            [
                "受講生のスキル合計を 8 グループに均等配分したい",
                "",
                "  Σᵢ skillᵢ xᵢ,g  =  AvgPoints  ∀g",
                "",
                "→ skill が整数で、AvgPoints が割り切れなければ",
                "  実行可能解が存在しない",
                "",
                "「絶対こうしたい」を絶対にすると、",
                "ちょっとした入力で全部壊れる",
            ],
            fill=RGBColor(0xFA, 0xE5, 0xD3), title_color=COLOR_ACCENT)
    add_box(s, 6.8, 1.7, 6.0, 5.0,
            "ソフト化: 偏差を最小化",
            [
                "  Σᵢ skillᵢ xᵢ,g  +  uᵍ  -  oᵍ  =  AvgPoints",
                "  uᵍ, oᵍ ≥ 0",
                "",
                "目的に偏差ペナルティを足す:",
                "  minimize   Σ_g ( uᵍ + oᵍ )",
                "",
                "★研究室文脈:",
                "  CIM/CAC のロス関数 = エネルギー + ペナルティ",
                "  「ペナルティ係数の重み」が実装上の自由度",
                "  ハード制約を直接書くより、ペナルティ化で",
                "  最適化問題に乗せ直すアプローチが多い",
            ],
            fill=RGBColor(0xE8, 0xF8, 0xF5),
            title_color=RGBColor(0x16, 0xA0, 0x85))
    add_footer(s, 13, 20)
    slides_meta.append("ソフト制約")

    # =======================================
    # Slide 14: 章扉 III
    # =======================================
    s = prs.slides.add_slide(blank)
    add_section_divider(s, "Part III", "研究室の研究との接続")
    slides_meta.append("Part III 扉")

    # =======================================
    # Slide 15: MaxCut を IP で書く
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "MaxCut を 0-1 IP として書く", "研究テーマを定式化の言葉で")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "G-set / G22 (N=2000) — MILP ソルバーの「典型問題」と同じ枠組みで書ける",
                size=14, color=COLOR_SUB)
    add_box(s, 0.6, 1.7, 6.0, 5.0,
            "(1) ナイーブな 0-1 IP",
            [
                "決定変数: xᵢ ∈ {0, 1}    頂点 i を集合 A か B かで分ける",
                "",
                "辺 (i, j) がカットされる ⇔ xᵢ ≠ xⱼ",
                "    ⇔  xᵢ + xⱼ - 2 xᵢ xⱼ  =  1",
                "",
                "目的:",
                "  maximize  Σ_(i,j)∈E  (xᵢ + xⱼ - 2 xᵢ xⱼ)",
                "",
                "→ 2 次項が出る  ⇒  QUBO 形式",
                "→ Cbc では解けない (Gurobi なら可)",
            ],
            fill=RGBColor(0xFA, 0xE5, 0xD3), title_color=COLOR_ACCENT)
    add_box(s, 6.8, 1.7, 6.0, 5.0,
            "(2) 線形化 IP (本物の MILP)",
            [
                "新変数 yᵢⱼ ∈ {0, 1} を導入 ── yᵢⱼ = (xᵢ XOR xⱼ)",
                "",
                "maximize  Σ_(i,j)∈E  yᵢⱼ",
                "",
                "subject to:",
                "  yᵢⱼ  ≤  xᵢ + xⱼ",
                "  yᵢⱼ  ≤  2 - xᵢ - xⱼ",
                "  yᵢⱼ  ≥  xᵢ - xⱼ",
                "  yᵢⱼ  ≥  xⱼ - xᵢ",
                "",
                "★MILP ソルバーで解ける形に変換可能",
                "  だが N=2000 で辺 19990 本では現実時間で解けない",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_footer(s, 15, 20)
    slides_meta.append("MaxCut IP")

    # =======================================
    # Slide 16: MILPソルバー vs CIM
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "MILP ソルバー vs CIM / CAC", "アプローチの住み分け")

    headers = ["", "MILP ソルバー (Gurobi/Cbc)", "CIM / CAC / SA / SB"]
    rows = [
        ["定式化", "0-1 IP / QUBO の線形化",
         "Ising ハミルトニアン  H = -Σ Jᵢⱼ σᵢ σⱼ"],
        ["決定変数", "xᵢ ∈ {0, 1}",
         "σᵢ ∈ {-1, +1}  (光振幅 cᵢ の符号)"],
        ["保証", "最適性ギャップを数値で出せる",
         "近似解。最適性証明なし"],
        ["小規模", "N ≤ 100 なら数秒で厳密最適解",
         "オーバーキル"],
        ["G-set 規模", "N=2000 では時間切れになる場合あり",
         "数百 ms で best_cut ≈ 13320 (G22)"],
        ["パラメータ", "ほぼなし (gap, time limit 程度)",
         "kappa, gamma, dP/round, 損失 dB …多数"],
        ["役割", "「正解はこれ」と保証する基準",
         "規模を上げたときの実用的な解法"],
    ]
    n_cols = 3
    n_rows = len(rows) + 1
    table_shape = s.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.4),
        Inches(12.3), Inches(5.4),
    )
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(5.0)
    table.columns[2].width = Inches(5.3)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = FONT_JA
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.name = FONT_JA
            run.font.size = Pt(11)
            run.font.color.rgb = COLOR_SUB
            if c == 0:
                run.font.bold = True
                run.font.color.rgb = COLOR_PRIMARY
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEC, 0xF2, 0xF9)
    add_textbox(s, 0.5, 6.9, 12.3, 0.3,
                "→ どちらか一方ではなく、「小規模は MILP で答合わせ、大規模は CIM」が現実的",
                size=11, bold=True, color=COLOR_ACCENT)
    add_footer(s, 16, 20)
    slides_meta.append("MILP vs CIM")

    # =======================================
    # Slide 17: Optuna も最適化問題
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Optuna も「最適化問題」だった", "連続・ブラックボックス最適化として")
    add_textbox(s, 0.6, 1.1, 12.0, 0.5,
                "BrainPad の分類でいうと「連続 × 非線形 × ブラックボックス目的関数」に該当",
                size=14, color=COLOR_SUB)
    add_box(s, 0.6, 1.7, 6.0, 5.0,
            "Optuna のお気持ちを LP/IP の言葉で",
            [
                "決定変数 (連続):",
                "  L, gamma, loss_dB, dP_per_round, coupling",
                "",
                "目的関数 (ブラックボックス):",
                "  f(params) = 20 試行平均の best_cut",
                "",
                "制約条件:",
                "  各パラの探索範囲 (suggest_float の min/max)",
                "",
                "→ 「数式で目的関数が書けない」点だけが LP と違う",
                "→ TPE 等のサンプラーで賢く点を打って探す",
            ],
            fill=RGBColor(0xFD, 0xF6, 0xE3), title_color=COLOR_ACCENT)
    add_box(s, 6.8, 1.7, 6.0, 5.0,
            "今週やった num_rounds スイープの位置づけ",
            [
                "外側ループ (人がやる): num_rounds ∈ {30, 300, 3k, 10k}",
                "内側ループ (Optuna):  5 パラを TPE で探索",
                "  → 計算予算 = (rounds × n_trials) ほぼ一定",
                "",
                "得られた知見:",
                "  ・30 rounds: 論文値 10337 → Optuna 12886 (+2549)",
                "  ・3000+ rounds: 既に飽和し改善幅 +20 程度",
                "",
                "= 数理最適化の枠組みで「最適化を最適化」している",
                "  (メタ最適化)",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_footer(s, 17, 20)
    slides_meta.append("Optuna")

    # =======================================
    # Slide 18: 今後研究室で活かす
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "今後 研究室でこう活かす", "学んだことの再翻訳")
    add_bullets(s, 0.6, 1.2, 12.0, 5.5, [
        ("自分の問題を「3 要素」で書く癖", 0),
        ("CIM 研究でも、評価関数を変えるときは『目的・変数・制約』に分解する", 1),
        ("ペナルティ係数の議論は「ソフト制約のチューニング」と同じ", 1),
        ("", 0),
        ("典型問題のカタログを持つ", 0),
        ("ナップサック / 施設配置 / 割当 / スケジューリングは応用範囲が広い", 1),
        ("MaxCut の派生(重み付き / k-cut / 制約付き)を IP として書けるようにする", 1),
        ("", 0),
        ("小規模ベンチマークで答え合わせをする", 0),
        ("N ≤ 100 の MaxCut なら Gurobi で厳密最適を求めて CIM と比較できる", 1),
        ("「13321 が本当の最適か」を MILP ソルバーで証明する逆方向の研究も可能", 1),
        ("", 0),
        ("メタ最適化を意識的に設計する", 0),
        ("Optuna 探索空間・予算配分・warm start は「最適化問題の設計問題」と捉える", 1),
        ("BrainPad の「ハード→ソフト制約」の発想で、Optuna の探索を頑健にできる", 1),
    ], size=13)
    add_footer(s, 18, 20)
    slides_meta.append("活かし方")

    # =======================================
    # Slide 19: まとめ
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "まとめ", "持ち帰る 3 つのメッセージ")
    msgs = [
        ("①", "数理最適化 = 目的関数 / 決定変数 / 制約 の 3 要素で世界を書き直す技術",
         "  自分の研究も「何を最大化したいか」を言語化する練習になる"),
        ("②", "LP / IP / MIP は「問題のタイプ」のラベル。MaxCut は 0-1 IP の典型",
         "  CIM はその「フルスクラッチ・問題固有」アプローチ"),
        ("③", "ソルバーは選ぶもの。汎用ソルバー / CIM / Optuna は同じ最適化の家族",
         "  問題規模と保証要件で住み分ければよい"),
    ]
    for i, (num, head, tail) in enumerate(msgs):
        top = 1.4 + i * 1.8
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.6), Inches(top), Inches(0.9), Inches(0.9),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.fill.background()
        add_textbox(s, 0.6, top + 0.18, 0.9, 0.6,
                    num, size=22, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_textbox(s, 1.7, top, 11.0, 0.5,
                    head, size=16, bold=True, color=COLOR_PRIMARY)
        add_textbox(s, 1.7, top + 0.55, 11.0, 0.5,
                    tail, size=12, color=COLOR_SUB)
    add_footer(s, 19, 20)
    slides_meta.append("まとめ")

    # =======================================
    # Slide 20: 参考文献 / 次に勉強すること
    # =======================================
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "参考文献 / 今後の学習リスト", "⑧ 今後学ぶと良いこと")
    add_box(s, 0.6, 1.2, 6.0, 5.5,
            "典型問題を学ぶ書籍",
            [
                "「Python ではじめる数理最適化」",
                "  (BrainPad DS 新卒指定書籍)",
                "  実務寄り、モデリング多数",
                "",
                "「あたらしい数理最適化」",
                "  典型問題のコツが詰まっている",
                "",
                "Williams 「Model Building in Mathematical Programming」",
                "  MIP モデリングの聖典",
                "",
                "梅谷俊治 「しっかり学ぶ数理最適化」",
                "  理論寄り (BrainPad 資料の引用元)",
            ],
            fill=RGBColor(0xEC, 0xF2, 0xF9), title_color=COLOR_PRIMARY)
    add_box(s, 6.8, 1.2, 6.0, 5.5,
            "次のステップ (研究との結合)",
            [
                "ヒューリスティクスを学ぶ",
                "  AtCoder Heuristic Contest (AHC)",
                "  焼きなまし・ビームサーチ etc.",
                "  → CIM/CAC を SA と並べる審美眼が育つ",
                "",
                "他ソルバーに触れる",
                "  OR-Tools CP-SAT (論理制約に強い)",
                "  Gurobi (有償だが学術ライセンス可)",
                "",
                "Web アプリ化",
                "  研究結果をデモにしたいときに必要",
                "  Streamlit + PuLP の組合せが簡単",
            ],
            fill=RGBColor(0xFD, 0xF6, 0xE3), title_color=COLOR_ACCENT)

    add_textbox(s, 0.6, 6.85, 12.0, 0.3,
                "原典: BrainPad『数理最適化研修』(社内公開資料, 2025) を研究室文脈で再構成 — 内容の責任は本発表者",
                size=10, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(s, 20, 20)
    slides_meta.append("参考文献")

    # ============================================================
    # 保存
    # ============================================================
    out_path = Path("docs/optimization_training_report.pptx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
    for i, name in enumerate(slides_meta, start=1):
        print(f"  {i:2d}. {name}")


if __name__ == "__main__":
    build()
